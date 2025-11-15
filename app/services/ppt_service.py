# app/services/ppt_service.py
import os
import json
import logging
import re
from typing import Dict, List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openai import OpenAI

# Config 不再需要，因为模板路径和输出路径都是通过参数传入的

logger = logging.getLogger(__name__)

# 初始化Qwen API客户端
QWEN_API_KEY = os.getenv("QWEN_API_KEY") or "sk-2633adfd6c8a43819bcedd895eb79c0f"
QWEN_BASE_URL = os.getenv("QWEN_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1")
QWEN_MODEL = os.getenv("QWEN_MODEL", "qwen-plus")

if not QWEN_API_KEY:
    logger.error("QWEN_API_KEY 未设置，PPT生成功能将无法正常工作")
else:
    logger.info(f"QWEN_API_KEY 已配置（长度: {len(QWEN_API_KEY)}）")

qwen_client = OpenAI(
    api_key=QWEN_API_KEY,
    base_url=QWEN_BASE_URL,
)

# PPT模板文件路径
_project_root = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
PPT_TEMPLATE_PATH_PPTX = os.path.join(_project_root, "个人竞聘_求职汇报模板.pptx")
PPT_TEMPLATE_PATH_POTX = os.path.join(_project_root, "个人竞聘_求职汇报模板.potx")

if os.path.exists(PPT_TEMPLATE_PATH_PPTX):
    PPT_TEMPLATE_PATH = PPT_TEMPLATE_PATH_PPTX
elif os.path.exists(PPT_TEMPLATE_PATH_POTX):
    PPT_TEMPLATE_PATH = PPT_TEMPLATE_PATH_POTX
else:
    PPT_TEMPLATE_PATH = None


def build_ppt_outline_prompt(resume_text: str, jd_text: str) -> str:
    """构建用于生成PPT大纲的prompt"""
    return f"""# 角色

你是一位专业的PPT内容规划专家。你的任务是根据简历和岗位JD，提取并整理信息，生成PPT大纲。

# 任务

从简历中提取信息，结合岗位JD的要求，生成以下结构的PPT大纲。所有信息必须严格来自简历，不能篡改或夸大事实。

# 输出格式

请只输出JSON格式，不要包含任何其他文字或代码块标记。

JSON格式：
{{
  "title": "姓名+个人简介（例如：张伟个人简介）",
  "name": "姓名（从简历中提取）",
  "1": "姓名、出生年月、电话、邮箱、籍贯、政治面貌（每项一行，格式：姓名：XXX）",
  "2": "学校、专业、主修课、就读时间（每项一行，格式：学校：XXX，专业：XXX，主修课：XXX（多个主修课用逗号分隔），就读时间：XXX）",
  "3": "学分绩点、奖学金（每项一行，格式：学分绩点：XXX；奖学金直接写名称，不要写"奖学金："前缀）",
  "4": "项目名称、时间（每个项目一行，格式：项目名称：XXX，时间：XXX）",
  "5": "负责内容（所有项目的负责内容合并，直接写内容，不要写"负责内容："前缀，整体内容不超过300字，用简洁的语言概括所有项目的负责内容，不要以省略号结尾）",
  "6": "项目成果（每个项目的成果，直接写成果，不要写"项目成果："前缀，每个项目一行）",
  "7": "实习与工作经历（每段经历格式：公司名称\\n职位\\n时间\\n工作内容（工作内容可以是多行，每行一个要点）。不同工作经历之间用一个空行分隔。如无实习或工作经历则输出：无）",
  "8": "个人技能（技术技能和软技能，优先选择与岗位JD相关的，按照以下格式整理：精通xx、xx\\n熟悉xx、xx\\n后端：xx、xx\\n前端：xx、xx\\n数据库：xx、xx\\n其他：xx、xx。如果某个分类没有内容，则省略该分类）",
  "9": "未来职业展望（基于简历和岗位JD，包含1年职业展望、3年职业展望等，格式：1年职业展望：XXX\\n3年职业展望：XXX。至少包含1年和3年展望，可以适当扩展）"
}}

# 重要要求

1. **事实准确性**：所有信息必须严格来自简历，不能篡改、夸大或虚构
2. **内容匹配**：在"个人技能"和"未来职业展望"部分，优先选择与岗位JD相关的技能和展望
3. **格式要求**：
   - 每个字段的内容用换行符分隔不同项
   - 如果某项信息在简历中不存在，输出"无"或省略该项
   - 不要包含占位符标记如[从简历中提取]
   - 负责内容部分整体不超过300字（所有项目的负责内容合并，用简洁的语言概括，不要以省略号结尾）
   - 技能部分必须按照指定格式整理（精通、熟悉、后端、前端、数据库、其他等分类）
   - 未来展望部分必须包含1年职业展望和3年职业展望，可以适当扩展
   - 实习与工作经历部分使用换行格式：公司名称\\n职位\\n时间\\n工作内容，不同工作经历之间用一个空行分隔
4. **输出要求**：只输出JSON，不要有任何前缀、后缀或代码块标记

# 输入信息

【简历内容】:
---
{resume_text}
---

【岗位JD】:
---
{jd_text}
---

# 输出

请直接输出JSON格式的PPT大纲。
"""


def generate_ppt_outline(resume_text: str, jd_text: str) -> Dict:
    """使用Qwen API生成PPT大纲"""
    prompt = build_ppt_outline_prompt(resume_text, jd_text)
    logger.info(f"开始生成PPT大纲，简历长度: {len(resume_text)} 字符，JD长度: {len(jd_text)} 字符")
    
    try:
        completion = qwen_client.chat.completions.create(
            model=QWEN_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "你是一位专业的PPT内容规划专家。你的任务是根据简历和岗位JD，提取并整理信息，生成PPT大纲。所有信息必须严格来自简历，不能篡改或夸大事实。"
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.2,
        )
        
        response = completion.choices[0].message.content.strip()
        
        if not response or len(response.strip()) < 10:
            logger.warning("模型返回内容过短")
            raise ValueError("模型返回内容过短")
            
    except Exception as e:
        logger.error(f"调用Qwen API失败: {e}")
        raise
    
    # 提取JSON
    json_str = response.strip()
    json_str = re.sub(r'^```json\s*', '', json_str, flags=re.MULTILINE)
    json_str = re.sub(r'^```\s*', '', json_str, flags=re.MULTILINE)
    json_str = re.sub(r'```\s*$', '', json_str, flags=re.MULTILINE)
    json_str = json_str.strip()
    
    # 查找JSON对象
    json_start = json_str.find('{')
    json_end = json_str.rfind('}')
    
    if json_start != -1 and json_end != -1 and json_end > json_start:
        json_str = json_str[json_start:json_end+1]
    
    try:
        outline = json.loads(json_str)
        logger.info("成功解析PPT大纲")
        return outline
    except json.JSONDecodeError as e:
        logger.error(f"JSON解析失败: {e}")
        logger.error(f"响应内容前500字符: {json_str[:500]}")
        raise


def find_textbox_by_content(slide, search_text: str):
    """
    根据文本内容查找文本框
    优先查找精确匹配的文本框（文本内容就是搜索文本）
    """
    # 先查找精确匹配的文本框（文本内容就是搜索文本）
    exact_matches = []
    # 再查找包含搜索文本的文本框
    partial_matches = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text:
            text = shape.text_frame.text.strip()
            if text == search_text:
                # 精确匹配
                exact_matches.append(shape)
            elif search_text in text:
                # 部分匹配
                partial_matches.append(shape)
    
    # 优先返回精确匹配
    if exact_matches:
        return exact_matches[0]
    elif partial_matches:
        return partial_matches[0]
    else:
        logger.warning(f"未找到文本框: '{search_text}'")
        return None


def fill_textbox(shape, content: str, preserve_format: bool = True):
    """
    填充文本框内容，保留原格式
    
    Args:
        shape: 文本框形状
        content: 要填充的内容
        preserve_format: 是否保留原格式（默认True）
    """
    if not shape or not hasattr(shape, 'text_frame'):
        return False
    
    try:
        text_frame = shape.text_frame
        
        # 如果保留格式，先保存原格式
        if preserve_format and len(text_frame.paragraphs) > 0:
            # 获取第一个段落的格式作为参考
            ref_paragraph = text_frame.paragraphs[0]
            ref_alignment = ref_paragraph.alignment if hasattr(ref_paragraph, 'alignment') else None
            ref_space_after = ref_paragraph.space_after if hasattr(ref_paragraph, 'space_after') else None
            ref_line_spacing = ref_paragraph.line_spacing if hasattr(ref_paragraph, 'line_spacing') else None
            
            # 保存字体属性 - 优先从runs中获取，如果没有runs则从段落font获取
            font_size = None
            font_name = None
            font_bold = None
            font_italic = None
            font_color = None
            
            # 优先从runs中获取格式（因为格式通常在runs中）
            if len(ref_paragraph.runs) > 0:
                ref_run = ref_paragraph.runs[0]
                ref_run_font = ref_run.font if hasattr(ref_run, 'font') else None
                if ref_run_font:
                    try:
                        font_size = ref_run_font.size
                        font_name = ref_run_font.name
                        font_bold = ref_run_font.bold
                        font_italic = ref_run_font.italic
                        if hasattr(ref_run_font, 'color') and ref_run_font.color:
                            font_color = ref_run_font.color.rgb if hasattr(ref_run_font.color, 'rgb') else None
                    except:
                        pass
            
            # 如果runs中没有格式，尝试从段落font获取
            if font_size is None:
                ref_font = ref_paragraph.font if hasattr(ref_paragraph, 'font') else None
                if ref_font:
                    try:
                        font_size = ref_font.size
                        font_name = ref_font.name
                        font_bold = ref_font.bold
                        font_italic = ref_font.italic
                        if hasattr(ref_font, 'color') and ref_font.color:
                            font_color = ref_font.color.rgb if hasattr(ref_font.color, 'rgb') else None
                    except:
                        pass
        else:
            # 使用默认格式
            ref_paragraph = None
            ref_font = None
            ref_alignment = None
            ref_space_after = None
            ref_line_spacing = None
            font_size = Pt(18)
            font_name = None
            font_bold = None
            font_italic = None
            font_color = None
        
        # 清空文本框
        text_frame.clear()
        text_frame.word_wrap = True
        
        # 如果内容是多行，按行分割（保留空行）
        lines = content.split('\n') if content else []
        # 不要strip和过滤空行，保留空行用于分隔
        processed_lines = []
        for line in lines:
            # 保留空行（用于分隔不同工作经历等）
            if not line.strip():
                processed_lines.append('')
            else:
                processed_lines.append(line.strip())
        
        lines = processed_lines
        # 检查是否有非空行
        if not any(line.strip() for line in lines):
            return False
        
        # 填充内容（处理空行）
        paragraph_count = 0
        for i, line in enumerate(lines):
            # 如果是空行，创建空段落
            if not line.strip():
                if paragraph_count == 0:
                    # 第一个段落不能为空，如果需要空行，创建新段落
                    p = text_frame.add_paragraph()
                    p.text = ''
                    paragraph_count += 1
                else:
                    # 创建空段落
                    p = text_frame.add_paragraph()
                    p.text = ''
                    paragraph_count += 1
                continue
            
            # 非空行
            if paragraph_count == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            paragraph_count += 1
            
            # 应用格式 - 需要应用到runs中
            if len(p.runs) > 0:
                run = p.runs[0]
                if hasattr(run, "font"):
                    if font_size:
                        try:
                            run.font.size = font_size
                        except:
                            pass
                    if font_name:
                        try:
                            run.font.name = font_name
                        except:
                            pass
                    if font_bold is not None:
                        try:
                            run.font.bold = font_bold
                        except:
                            pass
                    if font_italic is not None:
                        try:
                            run.font.italic = font_italic
                        except:
                            pass
                    if font_color:
                        try:
                            run.font.color.rgb = font_color
                        except:
                            pass
                    elif not preserve_format:
                        # 如果没有保留格式，使用默认黑色
                        try:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        except:
                            pass
            else:
                # 如果没有runs，尝试应用到段落font
                if hasattr(p, "font"):
                    if font_size:
                        try:
                            p.font.size = font_size
                        except:
                            pass
                    if font_name:
                        try:
                            p.font.name = font_name
                        except:
                            pass
                    if font_bold is not None:
                        try:
                            p.font.bold = font_bold
                        except:
                            pass
                    if font_italic is not None:
                        try:
                            p.font.italic = font_italic
                        except:
                            pass
                    if font_color:
                        try:
                            p.font.color.rgb = font_color
                        except:
                            pass
                    elif not preserve_format:
                        # 如果没有保留格式，使用默认黑色
                        try:
                            p.font.color.rgb = RGBColor(0, 0, 0)
                        except:
                            pass
            
            # 应用段落格式
            if ref_alignment is not None:
                try:
                    p.alignment = ref_alignment
                except:
                    p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.LEFT
            
            if ref_space_after is not None:
                try:
                    p.space_after = ref_space_after
                except:
                    p.space_after = Pt(8)
            else:
                p.space_after = Pt(8)
            
            p.level = 0
        
        return True
    except Exception as e:
        logger.error(f"填充文本框时出错: {e}", exc_info=True)
        return False


def create_ppt_from_template(outline: Dict, template_path: str, output_path: str) -> str:
    """
    使用模板文件创建PPT，按照指定位置填充内容
    """
    from datetime import datetime
    import shutil
    import tempfile
    
    # 检查模板文件
    if not template_path or not os.path.exists(template_path):
        logger.error(f"模板文件不存在: {template_path}")
        raise FileNotFoundError(f"模板文件不存在: {template_path}")
    
    # 处理.potx文件
    actual_template_path = template_path
    temp_pptx_path = None
    
    if template_path.lower().endswith('.potx'):
        try:
            temp_pptx_path = tempfile.mktemp(suffix='.pptx')
            shutil.copy2(template_path, temp_pptx_path)
            actual_template_path = temp_pptx_path
        except Exception as e:
            logger.warning(f"转换.potx文件失败: {e}")
            actual_template_path = template_path
    
    try:
        prs = Presentation(actual_template_path)
    except Exception as e:
        logger.error(f"加载模板文件失败: {e}")
        raise
    
    current_date = datetime.now().strftime("%Y年%m月%d日")
    name = outline.get('name', '候选人')
    title = outline.get('title', f'{name}个人简介')
    
    # 填充第一页：标题和汇报人
    if len(prs.slides) > 0:
        slide1 = prs.slides[0]
        
        # 填充标题
        title_box = find_textbox_by_content(slide1, "标题")
        if title_box:
            fill_textbox(title_box, title, preserve_format=True)
        
        # 填充汇报人
        reporter_box = find_textbox_by_content(slide1, "汇报人")
        if reporter_box:
            fill_textbox(reporter_box, name, preserve_format=True)
        
        # 填充日期
        date_box = find_textbox_by_content(slide1, "汇报日期")
        if date_box:
            fill_textbox(date_box, f"汇报日期：{current_date}", preserve_format=True)
    
    # 填充第三页：1（个人基本信息）
    if len(prs.slides) >= 3:
        slide3 = prs.slides[2]  # 索引2是第三页
        box1 = find_textbox_by_content(slide3, "1")
        if box1:
            content1 = outline.get('1', '')
            if content1:
                fill_textbox(box1, content1, preserve_format=True)
    
    # 填充第四页：2和3（教育经历）
    if len(prs.slides) >= 4:
        slide4 = prs.slides[3]  # 索引3是第四页
        
        # 填充2（学校、专业、主修课、就读时间）
        box2 = find_textbox_by_content(slide4, "2")
        if box2:
            content2 = outline.get('2', '')
            if content2:
                # 清理主修课前缀
                content2 = re.sub(r'主修课：', '', content2)
                fill_textbox(box2, content2, preserve_format=True)
        
        # 填充3（学分绩点、奖学金）
        box3 = find_textbox_by_content(slide4, "3")
        if box3:
            content3 = outline.get('3', '')
            if content3:
                # 清理奖学金前缀
                content3 = re.sub(r'奖学金：', '', content3)
                fill_textbox(box3, content3, preserve_format=True)
    
    # 填充第五页：4、5、6（项目经历）
    if len(prs.slides) >= 5:
        slide5 = prs.slides[4]  # 索引4是第五页
        
        # 填充4（项目名称、时间）
        box4 = find_textbox_by_content(slide5, "4")
        if box4:
            content4 = outline.get('4', '')
            if content4:
                fill_textbox(box4, content4, preserve_format=True)
        
        # 填充5（负责内容）
        box5 = find_textbox_by_content(slide5, "5")
        if box5:
            content5 = outline.get('5', '')
            if content5:
                # 清理负责内容前缀
                content5 = re.sub(r'负责内容：', '', content5)
                content5 = re.sub(r'^负责内容\s*[：:]\s*', '', content5, flags=re.MULTILINE)
                # 移除换行符，合并为一段文字
                content5 = re.sub(r'\n+', ' ', content5).strip()
                
                # 移除末尾的省略号（如果API返回的内容以...结尾）
                # 移除末尾的多个点号（...或....等），包括前面可能有中文字符的情况
                # 先移除"多..."这样的模式
                content5 = re.sub(r'多\.{2,}\s*$', '', content5).strip()
                # 再移除末尾的多个点号
                content5 = re.sub(r'\.{2,}\s*$', '', content5).strip()
                
                # 不再在代码中截断，因为prompt已经限制了长度
                # 如果内容确实超过300字，记录警告但不截断（让用户知道需要调整prompt）
                char_count = 0
                for char in content5:
                    if '\u4e00' <= char <= '\u9fff':
                        char_count += 2  # 中文字符
                    else:
                        char_count += 1  # 其他字符
                
                if char_count > 300:
                    logger.warning(f"负责内容超过300字（实际{char_count}字），但已按prompt要求生成，不进行截断")
                
                fill_textbox(box5, content5, preserve_format=True)
        
        # 填充6（项目成果）
        box6 = find_textbox_by_content(slide5, "6")
        if box6:
            content6 = outline.get('6', '')
            if content6:
                # 清理项目成果前缀
                content6 = re.sub(r'项目成果：', '', content6)
                content6 = re.sub(r'^项目成果\s*[：:]\s*', '', content6, flags=re.MULTILINE)
                fill_textbox(box6, content6, preserve_format=True)
    
    # 填充第六页：7（实习与工作经历）
    if len(prs.slides) >= 6:
        slide6 = prs.slides[5]  # 索引5是第六页
        box7 = find_textbox_by_content(slide6, "7")
        if box7:
            content7 = outline.get('7', '无')
            if content7 and content7 != '无':
                # 处理实习与工作经历格式
                # 格式：公司名称\n职位\n时间\n工作内容
                # 不同工作经历之间用一个空行分隔
                
                # 清理内容中的前缀
                content7 = re.sub(r'公司名称[：:]\s*', '', content7)
                content7 = re.sub(r'职位[：:]\s*', '', content7)
                content7 = re.sub(r'时间[：:]\s*', '', content7)
                content7 = re.sub(r'工作内容[：:]\s*', '', content7)
                
                # 按行分割（保留所有行，包括空行）
                all_lines = [line.rstrip() for line in content7.split('\n')]
                # 移除空行，重新处理
                lines = [line.strip() for line in all_lines if line.strip()]
                
                experiences = []
                i = 0
                
                # 简化的解析逻辑：按模式识别工作经历
                # 模式：公司名称 -> 职位 -> 时间 -> 工作内容（多行）-> 下一个公司名称
                while i < len(lines):
                    # 查找公司名称（不包含时间格式的短行）
                    if i < len(lines) and not re.search(r'\d{4}.*[年月]', lines[i]) and not re.search(r'\d{4}.*-.*\d{4}', lines[i]):
                        company = lines[i]
                        position = None
                        time = None
                        work_content_lines = []
                        
                        i += 1
                        
                        # 查找职位（下一个短行，不包含时间）
                        if i < len(lines) and len(lines[i]) < 30 and not re.search(r'\d{4}', lines[i]):
                            position = lines[i]
                            i += 1
                        
                        # 查找时间（包含年月日）
                        if i < len(lines) and (re.search(r'\d{4}.*[年月]', lines[i]) or re.search(r'\d{4}.*-.*\d{4}', lines[i])):
                            time = lines[i]
                            i += 1
                        
                        # 收集工作内容（直到下一个公司名称）
                        while i < len(lines):
                            line = lines[i]
                            
                            # 检查是否是下一个公司名称（短行，且接下来的行可能是职位或时间）
                            if len(line) < 50 and not re.search(r'\d{4}', line):
                                # 检查接下来的2-3行
                                look_ahead_count = 0
                                has_next_company_pattern = False
                                
                                # 检查下一行是否是职位或时间
                                if i + 1 < len(lines):
                                    next_line = lines[i + 1]
                                    # 如果是时间，说明这是下一个经历的开始
                                    if re.search(r'\d{4}.*[年月]', next_line) or re.search(r'\d{4}.*-.*\d{4}', next_line):
                                        has_next_company_pattern = True
                                    # 如果是职位（短行），再检查下下行是否是时间
                                    elif len(next_line) < 30 and i + 2 < len(lines):
                                        next_next_line = lines[i + 2]
                                        if re.search(r'\d{4}.*[年月]', next_next_line) or re.search(r'\d{4}.*-.*\d{4}', next_next_line):
                                            has_next_company_pattern = True
                                
                                if has_next_company_pattern:
                                    # 这是下一个经历的开始
                                    break
                            
                            # 否则，这是工作内容
                            work_content_lines.append(line)
                            i += 1
                        
                        # 保存经历
                        if company and time:
                            # 工作内容保持换行（如果有多行）
                            # 如果工作内容是用分号或逗号分隔的，尝试转换为换行格式
                            work_content = '\n'.join(work_content_lines).strip()
                            # 如果工作内容很长且没有换行，可能是用分号或逗号分隔的
                            if work_content and '\n' not in work_content and (';' in work_content or '；' in work_content):
                                # 尝试按分号分割
                                work_items = re.split(r'[;；]', work_content)
                                work_items = [item.strip() for item in work_items if item.strip()]
                                work_content = '\n'.join(work_items)
                            experiences.append([company, position or '', time, work_content])
                    else:
                        i += 1
                
                # 格式化输出：每个经历格式为：公司名称\n职位\n时间\n工作内容
                # 不同工作经历之间用一个空行分隔
                formatted_lines = []
                for idx, exp in enumerate(experiences):
                    company = exp[0] if len(exp) > 0 else ''
                    position = exp[1] if len(exp) > 1 else ''
                    time = exp[2] if len(exp) > 2 else ''
                    work_content = exp[3] if len(exp) > 3 else ''
                    
                    if company:
                        formatted_lines.append(company)
                        if position:
                            formatted_lines.append(position)
                        if time:
                            formatted_lines.append(time)
                        if work_content:
                            # 如果工作内容包含换行，保持换行
                            if '\n' in work_content:
                                work_lines = work_content.split('\n')
                                formatted_lines.extend(work_lines)
                            else:
                                formatted_lines.append(work_content)
                        
                        # 如果不是最后一个经历，添加空行
                        if idx < len(experiences) - 1:
                            formatted_lines.append('')
                
                if formatted_lines:
                    content7 = '\n'.join(formatted_lines)
                # 如果没有处理成功，保持原格式
                
                fill_textbox(box7, content7, preserve_format=True)
    
    # 填充第七页：8（个人技能）
    if len(prs.slides) >= 7:
        slide7 = prs.slides[6]  # 索引6是第七页
        box8 = find_textbox_by_content(slide7, "8")
        if box8:
            content8 = outline.get('8', '')
            if content8:
                # 技能部分已经按照格式整理（精通、熟悉、后端、前端、数据库、其他等分类）
                # 不需要清理前缀，直接填充
                fill_textbox(box8, content8, preserve_format=True)
    
    # 填充第八页：9（未来职业展望）
    if len(prs.slides) >= 8:
        slide8 = prs.slides[7]  # 索引7是第八页
        box9 = find_textbox_by_content(slide8, "9")
        if box9:
            content9 = outline.get('9', '')
            if content9:
                # 未来展望部分已经按照格式整理（1年职业展望、3年职业展望等）
                # 不需要清理前缀，直接填充
                fill_textbox(box9, content9, preserve_format=True)
    
    # 填充第九页：汇报人
    if len(prs.slides) >= 9:
        slide9 = prs.slides[8]  # 索引8是第九页
        reporter_box = find_textbox_by_content(slide9, "汇报人")
        if reporter_box:
            fill_textbox(reporter_box, name, preserve_format=True)
        
        # 填充日期
        date_box = find_textbox_by_content(slide9, "汇报日期")
        if date_box:
            fill_textbox(date_box, f"汇报日期：{current_date}", preserve_format=True)
    
    # 保存PPT
    prs.save(output_path)
    logger.info(f"PPT已保存到: {output_path}")
    
    # 清理临时文件
    if temp_pptx_path and os.path.exists(temp_pptx_path):
        try:
            os.remove(temp_pptx_path)
        except:
            pass
    
    return output_path


def generate_self_intro_ppt(resume_text: str, jd_text: str, output_dir: str) -> Tuple[str, Optional[str]]:
    """
    生成自我介绍PPT
    
    步骤：
    1. 使用qwenAPI从简历和JD中提取信息，生成PPT大纲
    2. 使用模板文件填充数据到指定位置
    
    Returns:
        (ppt_path, error_message)
    """
    try:
        # 检查模板文件是否存在
        if not PPT_TEMPLATE_PATH or not os.path.exists(PPT_TEMPLATE_PATH):
            error_msg = f"PPT模板文件不存在: {PPT_TEMPLATE_PATH}"
            logger.error(error_msg)
            return "", error_msg
        
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
        
        # 步骤1: 生成PPT大纲
        outline = generate_ppt_outline(resume_text, jd_text)
        
        # 生成输出文件名
        import uuid
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        file_id = f"{timestamp}-{uuid.uuid4().hex[:8]}"
        ppt_filename = f"self_intro_{file_id}.pptx"
        ppt_path = os.path.join(output_dir, ppt_filename)
        
        # 步骤2: 使用模板文件填充内容
        create_ppt_from_template(outline, PPT_TEMPLATE_PATH, ppt_path)
        
        return ppt_path, None
    except Exception as e:
        logger.error(f"生成PPT时出错: {e}", exc_info=True)
        return "", str(e)
